"""
Document parser for the RFP Intelligence System.

Takes a file path, returns a list of LangChain Document objects with
clean text and useful metadata. Supports PDF, DOCX, PPTX, and MD.

The parser does NOT chunk. Chunking is a separate step (src/ingestion/chunker.py)
because the two change for different reasons.

Tables matter for RFPs — requirements, SLAs, and pricing usually live in
tables, not prose. So every format path extracts tabular content and renders
it as Markdown pipe tables inline with the surrounding text, where it stays
readable to both the embedder and the LLM.
"""
import re
import unicodedata
from pathlib import Path
from typing import Iterator, Literal
from uuid import uuid4

from langchain_core.documents import Document
from langchain_community.document_loaders import PyMuPDFLoader

import docx
from docx.document import Document as _DocxDocument
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table as _DocxTable
from docx.text.paragraph import Paragraph as _DocxParagraph
from pptx import Presentation


SourceType = Literal["company_kb", "incoming_rfp"]

# Page header/footer markers, e.g. "P a g e 13 | 22" or "Page 13 of 22".
_PAGE_MARKER = re.compile(r"^\s*p\s*a\s*g\s*e\b.*$", re.IGNORECASE)
# A TOC entry: text, a run of dot leaders, then a trailing page number.
_TOC_LINE = re.compile(r"^.*\.{3,}\s*\d+\s*$")
# Three or more dots (optionally space-separated) used as TOC leaders.
_DOT_LEADERS = re.compile(r"(?:\.\s?){3,}")
# A word split across a line break with a hyphen: "Deli-\nverables".
_HYPHEN_BREAK = re.compile(r"(\w)-\n(\w)")


def normalize_text(text: str) -> str:
    """
    Clean raw extracted text before chunking/embedding.

    Removes common PDF artifacts — running page-number headers/footers,
    table-of-contents dot leaders, and words hyphenated across line breaks —
    and collapses excess whitespace. Conservative by design: it strips clear
    noise only and leaves substantive text (including Markdown tables) intact.
    """
    if not text:
        return text

    text = unicodedata.normalize("NFKC", text)
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = _HYPHEN_BREAK.sub(r"\1\2", text)

    kept: list[str] = []
    for line in text.split("\n"):
        # Drop running page markers and pure TOC leader lines outright.
        if _PAGE_MARKER.match(line) or _TOC_LINE.match(line):
            continue
        line = _DOT_LEADERS.sub(" ", line)          # stray dot leaders mid-line
        line = re.sub(r"[ \t]+", " ", line).strip()  # collapse runs of spaces
        kept.append(line)

    cleaned = "\n".join(kept)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)  # squeeze blank-line runs
    return cleaned.strip()


def parse_file(file_path: str | Path, source_type: SourceType) -> list[Document]:
    """
    Parse a single file into LangChain Documents.

    One Document per page (PDF), per slide (PPTX), per file (DOCX, MD).
    Tables are extracted and rendered as Markdown inline with the text.
    Metadata is attached to every Document so it survives chunking.

    Args:
        file_path: Path to the file.
        source_type: Where this document came from — used at retrieval time
                     to filter by company knowledge vs incoming RFP.

    Returns:
        A list of Document objects, each with .page_content and .metadata.

    Raises:
        ValueError: If the file extension is unsupported.
        FileNotFoundError: If the file doesn't exist.
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"No such file: {path}")

    suffix = path.suffix.lower()
    doc_id = str(uuid4())

    if suffix == ".pdf":
        documents = _parse_pdf(path)
    elif suffix == ".docx":
        documents = _parse_docx(path)
    elif suffix == ".pptx":
        documents = _parse_pptx(path)
    elif suffix in (".md", ".txt"):
        documents = _parse_text(path)
    else:
        raise ValueError(
            f"Unsupported file type: {suffix}. "
            f"Supported: .pdf, .docx, .pptx, .md, .txt"
        )

    # Normalize text (strip PDF artifacts, fix hyphenation, collapse whitespace),
    # then drop any documents that ended up empty (e.g. a blank slide or page).
    for doc in documents:
        doc.page_content = normalize_text(doc.page_content)
    documents = [d for d in documents if d.page_content.strip()]

    # Attach common metadata to every Document
    for doc in documents:
        doc.metadata.update({
            "doc_id": doc_id,
            "source": path.name,
            "source_type": source_type,
        })

    return documents


# ---------------------------------------------------------------------------
# Table helpers
# ---------------------------------------------------------------------------

def _rows_to_markdown(rows: list[list[str]]) -> str:
    """
    Render a table (list of rows, each a list of cell strings) as a Markdown
    pipe table. The first row is treated as the header. Returns "" for an
    empty table so callers can skip it.
    """
    cleaned = [
        [(cell or "").replace("\n", " ").replace("|", "\\|").strip() for cell in row]
        for row in rows
        if any((cell or "").strip() for cell in row)
    ]
    if not cleaned:
        return ""

    width = max(len(r) for r in cleaned)
    cleaned = [r + [""] * (width - len(r)) for r in cleaned]  # pad ragged rows

    header = cleaned[0]
    body = cleaned[1:]

    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join(["---"] * width) + " |"]
    lines += ["| " + " | ".join(r) + " |" for r in body]
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def _parse_pdf(path: Path) -> list[Document]:
    """
    One Document per page. PyMuPDF preserves layout reasonably well and already
    sets 'page' in metadata. Any tables detected on a page are appended to that
    page's text as Markdown.
    """
    loader = PyMuPDFLoader(str(path))
    docs = loader.load()

    tables_by_page = _pdf_tables_by_page(path)
    for doc in docs:
        page = doc.metadata.get("page")
        table_md = tables_by_page.get(page)
        if table_md:
            doc.page_content = f"{doc.page_content}\n\n{table_md}"
    return docs


def _pdf_tables_by_page(path: Path) -> dict[int, str]:
    """
    Internal: detect tables per page using PyMuPDF and return a mapping of
    0-indexed page number -> Markdown for all tables on that page.

    Table detection is best-effort: a failure on any single page is swallowed
    so it never breaks ingestion of the page text.
    """
    import fitz  # PyMuPDF

    out: dict[int, str] = {}
    with fitz.open(str(path)) as pdf:
        for page_index, page in enumerate(pdf):
            try:
                found = page.find_tables()
            except Exception:
                continue
            blocks = []
            for table in found.tables:
                md = _rows_to_markdown(table.extract())
                if md:
                    blocks.append(md)
            if blocks:
                out[page_index] = "\n\n".join(blocks)
    return out


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def _iter_block_items(document: _DocxDocument) -> Iterator[_DocxParagraph | _DocxTable]:
    """
    Yield each paragraph and table in document order.

    python-docx exposes paragraphs and tables as separate flat lists, which
    loses their interleaving. Walking the body's XML children preserves reading
    order so a requirements table stays next to its heading.
    """
    body = document.element.body
    for child in body.iterchildren():
        if isinstance(child, CT_P):
            yield _DocxParagraph(child, document)
        elif isinstance(child, CT_Tbl):
            yield _DocxTable(child, document)


def _parse_docx(path: Path) -> list[Document]:
    """
    One Document per file. Walks paragraphs and tables in reading order so
    table content (which the old paragraph-only pass dropped entirely) is
    captured as Markdown.
    """
    document = docx.Document(str(path))
    parts: list[str] = []
    for block in _iter_block_items(document):
        if isinstance(block, _DocxParagraph):
            text = block.text.strip()
            if text:
                parts.append(text)
        else:  # _DocxTable
            rows = [[cell.text for cell in row.cells] for row in block.rows]
            md = _rows_to_markdown(rows)
            if md:
                parts.append(md)

    full_text = "\n\n".join(parts)
    return [Document(page_content=full_text, metadata={})]


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def _parse_pptx(path: Path) -> list[Document]:
    """One Document per slide. Pulls text and tables from every shape."""
    presentation = Presentation(str(path))
    documents = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_table:
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                md = _rows_to_markdown(rows)
                if md:
                    slide_parts.append(md)
            elif shape.has_text_frame and shape.text.strip():
                slide_parts.append(shape.text)
        if slide_parts:
            documents.append(Document(
                page_content="\n\n".join(slide_parts),
                metadata={"slide": slide_number},
            ))
    return documents


def _parse_text(path: Path) -> list[Document]:
    """One Document per file for plain text and Markdown."""
    text = path.read_text(encoding="utf-8", errors="replace")
    return [Document(page_content=text, metadata={})]
